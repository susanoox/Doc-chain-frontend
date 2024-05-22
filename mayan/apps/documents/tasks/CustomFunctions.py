
############################# custom Imports ##########################################
import requests
import json, cv2, io, hashlib
# import language_tool_python
from requests.exceptions import ConnectionError, Timeout
from requests.exceptions import JSONDecodeError
import logging, nltk
import pandas as pd

# Define a logger
logger = logging.getLogger(__name__)


########################################################################################
from pdf2image import convert_from_path
from PIL import Image
import pytesseract, hashlib
import pdfplumber
import PyPDF2, numpy as np
from docx import Document as worddoc
from pptx import Presentation
import subprocess

########################################################################################
from ..models.document_models import Document 
############################ URL Variables #############################################

NameSpaceList = ["dev", "test", "demo", "client"]
NameSpace = NameSpaceList[0]
BlockUrl = 'http://54.66.31.125:3001/filehash'
Ocrurl = 'http://3.107.59.224:8080/v2/ocr'
url_BOT = 'http://3.107.59.224:8080/v2/upload'
SummaryUrl = "http://3.107.59.224:8080/v2/summary" 
RequestTimeOut = 1200
text_content = ""

################################### Upload Functions ###################################


#################################### Custom Functions ##################################

def Bot_Content_Upload(data_for_BOT):
    try:
        print(data_for_BOT, url_BOT)
        response = requests.post(url_BOT, json=data_for_BOT, headers={'Content-Type': 'application/json'}, timeout=RequestTimeOut)
        response.raise_for_status()  # Raise an exception for 4xx and 5xx status codes
        if response.status_code == 200:
            print("File uploaded for Bot")
        else:
            print("Upload failed for bot")
    except ConnectionError as ce:
        print("Connection Error:", ce)
        # Handle connection errors here, such as retrying the request
    except Timeout as te:
        print("Timeout Error:", te)
        # Handle timeout errors here
    except requests.exceptions.RequestException as e:
        print("Request Exception:", e)
        # Handle other types of request exceptions here
        

            
def upload_to_blockchain(file_content, file_id):
    try:
        hash_value = hashlib.md5(file_content).hexdigest()
        url_BC = BlockUrl
        data_for_BC = {
            "fileHash": hash_value,
            "fileId": str(file_id)
        }
        print(data_for_BC)
        json_data = json.dumps(data_for_BC)
        response = requests.post(url_BC, data=json_data, headers={'Content-Type': 'application/json'}, timeout=RequestTimeOut)
        if response.status_code == 200:
            print("File uploaded for BlockChain")
        else:
            response = requests.post(url_BC, data=json_data, headers={'Content-Type': 'application/json'}, timeout=RequestTimeOut)
            if response.status_code == 200:
                print("File uploaded for BlockChain")
            else:
                print("Upload failed for BlockChain")
    except Exception as e:
        print("Error uploading to blockchain:", e)

def UploadSummary(payload):  
    try:
        print("SummaryUrl", SummaryUrl)
        response = requests.post(SummaryUrl, json=payload, timeout=RequestTimeOut)
        print("payload",payload, response)
        
        response_json = response.json()
        return response_json.get('response')
    except Exception as e:
        print("Error upload summary:", e)
    
        

#--------------------------------------- Tools -----------------------------------------------


def calculate_grammar_percentage(text):
    # Tokenize the text into sentences
    # sentences = nltk.sent_tokenize(text)

    # # Initialize a counter for correct sentences
    # correct_sentences = 0

    # # Iterate through each sentence
    # for sentence in sentences:
    #     # Tokenize the sentence into words
    #     words = nltk.word_tokenize(sentence)
        
    #     # Perform part-of-speech tagging
    #     pos_tags = nltk.pos_tag(words)
        
    #     # Count the number of verbs in the sentence
    #     num_verbs = sum(1 for _, tag in pos_tags if tag.startswith('VB'))

    #     # If the number of verbs is greater than 0, consider the sentence correct
    #     if num_verbs > 0:
    #         correct_sentences += 1

    # # Calculate the percentage of correct sentences
    # try:
    #     grammar_percentage = (correct_sentences / len(sentences)) * 100
    # except ZeroDivisionError:
    #     return 0
    # print("grammar_percentage : ",grammar_percentage)

    # return grammar_percentage
    return 0

    

# python -m nltk.downloader averaged_perceptron_tagger
# python -m nltk.downloader punkt


# def calculate_grammar_percentage(text):
#         # Initialize LanguageTool outside the function
#         global language_tool
#         if 'language_tool' not in globals():
#             language_tool = language_tool_python.LanguageTool('en-US')
#         # Split text once to avoid multiple split operations
#         words = text.split()
#         total_words = len(words)
#         # Batch process the text in chunks of 100 words
#         batch_size = 100
#         error_count = 0
#         for i in range(0, total_words, batch_size):
#             chunk = ' '.join(words[i:i+batch_size])
#             matches = language_tool.check(chunk)
#             error_count += len(matches)
#         correct_words = total_words - error_count
#         try:
#             grammar_percentage = (correct_words / total_words) * 100
#         except ZeroDivisionError:
#             return 0
#         print("grammar_percentage:", str(grammar_percentage))
#         return grammar_percentage
    
def extract_text_from_image_google(content):
    files = {'file': content}
    
    # Make a POST request to the OCR endpoint
    response = requests.post(Ocrurl, files=files)
    
    print(response)  # Print response status code for debugging
    
    try:
        # Attempt to parse the JSON response
        json_response = response.json()
        
        # Check if the request was successful (status code 200)
        if response.status_code == 200:
            # Extract the text from the response JSON
            text = json_response.get("response").get("ocr_text")
            return text
        else:
            print("Error:", response.status_code)
            return None
    except JSONDecodeError:
        print("Failed to decode JSON response")
        return None
    
##################################### End of Functions ##################################################

image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp')
ppt_extensions = ('.pptx')

################################### Function for doc format #########################################
def extract_text_from_doc(file_path):
    try:
        result = subprocess.run(["antiword", file_path], capture_output=True, text=True)
        return result.stdout
    except Exception as e:
        print(f"Error extracting text from .doc file: {e}")
        return None
    
def readFile(Data:Document):
    global text_content
    document_file = Data.file_latest
    print('document_file', document_file)
    ReadContent = ""
    #-------------------------------------------- Extract From Text Images ---------------------------------------------------------
    if str(document_file).lower().endswith(image_extensions):
            with document_file.file.open('rb') as img_file:
                content_text = extract_text_from_image_google(img_file.read())
                print(content_text)
                if content_text is not None:
                    ReadContent = content_text
            return ReadContent
    #----------------------------------- Extract From file  -----------------------------------------------------------------------------
    elif str(document_file).lower().endswith('.txt'):
        temp_content = ""
        print('document_file', document_file.file)
        with document_file.file.open('r') as file_handle:
            return file_handle.read()
    #----------------------------------- Extract From Excel File  -----------------------------------------------------------------------------
        
    # elif str(document_file).lower().endswith(('.xls', '.xlsx', '.xlsm', '.xlsb')):
    #     print("working xl..!")
    #     excel_data = pd.read_excel(document_file.file)
    #     print(excel_data)
    #     return str(excel_data)
    elif str(document_file).lower().endswith(('.xls', '.xlsx', '.xlsm', '.xlsb')):
        print("working xl..!")
        excel_data = pd.read_excel(document_file.file, sheet_name=None)  # Read all sheets
        content = ""
        for sheet_name, sheet_data in excel_data.items():
            content += f"Sheet: {sheet_name}\n"
            content += sheet_data.to_string(index=False, header=True)  # Convert to string
            content += "\n\n"
        print(content)
        return content
    #----------------------------------- Extract From CSV Document -----------------------------------------------------------------------------
    
    elif str(document_file).lower().endswith('.csv'):
        print("working csv..!")
        csv_data = pd.read_csv(document_file.file)  # Read CSV file
        print(csv_data)
        return csv_data.to_string(index=False, header=True)
    
    #----------------------------------- Extract From Word Document -----------------------------------------------------------------------------
    elif str(document_file).lower().endswith(('.docs', '.docx')):
        temp_content = ""
        print('document_file', document_file.file)
        with document_file.file.open('rb') as file_handle:
            doc = worddoc(file_handle)
            for paragraph in doc.paragraphs:
                temp_content += paragraph.text + '\n'
        print(temp_content)
        return temp_content
    #----------------------------------- Extract From Pdf -----------------------------------------------------------------------------
    elif str(document_file).lower().endswith('.pdf'):
        print("working doc...!")
        temp_content = ""
        print('document_file', document_file.file)

        with pdfplumber.open(document_file.file.path) as pdf:
            for page in pdf.pages:
                temp_content = temp_content + '\n' + page.extract_text(layout=True).strip()

        if (98 > calculate_grammar_percentage(temp_content)):
            print("Running inside of image loop")
            try:
                global text_content
                images = convert_from_path(document_file.file.path)
                for i, image in enumerate(images):
                    print("image", i)
                    try:
                        text = pytesseract.image_to_string(image)
                        if 80 > calculate_grammar_percentage(text) and len(temp_content.replace(" ", "")) < 40:
                            rgb_image = image.convert("RGB")
                            image_bytes = io.BytesIO()
                            rgb_image.save(image_bytes, format="JPEG")  # You can change the format if needed
                            image_bytes = image_bytes.getvalue()
                            content_text = extract_text_from_image_google(image_bytes)
                            text_content = text_content + content_text
                        else:
                            text_content = text_content + text
                    except Exception as e:
                        print("error while reading image",e)
                content = text_content
                return content
            except Exception as pdf_processing_error:
                    print(f"Error processing PDF: {pdf_processing_error}")
                    return None
        else:
            return temp_content
    # elif str(document_file).lower().endswith(('.ppt', '.pptx')):
    #     print("working ppt...!")
    #     temp_content = ""
    #     with document_file.file.open('rb') as file_handle:
    #         presentation = Presentation(file_handle)
    #         for slide in presentation.slides:
    #             for shape in slide.shapes:
    #                 if hasattr(shape, "text"):
    #                     temp_content += shape.text + '\n'
    #     print(temp_content)
    #     return temp_content    
    elif str(document_file).lower().endswith(ppt_extensions):
        print("working ppt...!")
        temp_content = ""
        try:
            with document_file.file.open('rb') as file_handle:
                presentation = Presentation(file_handle)
                for slide_num, slide in enumerate(presentation.slides, start=1):
                    temp_content += f"Slide {slide_num}:\n"
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                temp_content += paragraph.text + '\n'
                    temp_content += '\n'
            print(temp_content)
            return temp_content
        except Exception as e:
            print(f"Error processing PowerPoint file: {e}")
            return None
    elif str(document_file).lower().endswith('.doc'):
        print('working doc...!')
        try:
            text_content = extract_text_from_doc(document_file.file.path)
            if text_content is not None:
                return text_content
            else:
                return "Error: Unable to extract text from .doc file"
        except Exception as e:
            print(f"Error reading .doc file: {e}")
            return None
    else:
        return None
