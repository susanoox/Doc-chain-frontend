# pip install python-pptx
from pptx import Presentation

def read_pptx(file_path):
    # Load the presentation
    presentation = Presentation(file_path)
    
    # Iterate through slides
    for slide_number, slide in enumerate(presentation.slides):
        print(f"Slide {slide_number + 1}")
        
        # Iterate through the shapes in each slide
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)

if __name__ == "__main__":
    # Path to your PowerPoint file
    pptx_file = "/home/nagipragalathan/Test/mayan-edms-4.6.3/Backup/pptexamples.ppt"
    
    read_pptx(pptx_file)
